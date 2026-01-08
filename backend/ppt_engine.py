import math 
import os
import uuid
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from models import PresentationData

# === 1. 辅助函数 ===

def get_image_stream(query):
    # 1. 设置请求头（防止被网站拦截）
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
    }
    
    # 2. 尝试使用 Pollinations AI (生成图)
    # 将 query 中的空格替换为 %20
    safe_query = query.replace(" ", "%20")
    # 增加 nologo=true 去水印，设置宽高
    url = f"https://image.pollinations.ai/prompt/{safe_query}?width=1280&height=720&nologo=true"
    
    print(f"   ⬇️ [Image] 正在下载图片: {query}...")

    try:
        response = requests.get(url, headers=headers, timeout=15) # 超时稍微给长一点点
        # 3. 检查状态码，只有 200 才算成功
        if response.status_code == 200 and len(response.content) > 0:
            return BytesIO(response.content)
        else:
            print(f"   ⚠️ AI绘图失败 (Code: {response.status_code})，准备切换备用源...")
            
    except Exception as e:
        print(f"   ⚠️ AI绘图连接报错: {e}")

    # --- 4. 兜底方案 (如果上面失败了，用随机图) ---
    print("   🔄 尝试使用备用图源 (Picsum)...")
    try:
        # Picsum 是一个非常稳定的随机图源
        backup_url = "https://picsum.photos/1280/720"
        backup_resp = requests.get(backup_url, headers=headers, timeout=10)
        if backup_resp.status_code == 200:
            return BytesIO(backup_resp.content)
    except Exception as e:
        print(f"   ❌ 备用图源也失败了: {e}")

    # 5. 实在不行返回 None，渲染引擎里会跳过插图逻辑，防止程序崩溃
    return None

def auto_fit_text(text_frame, content_list: list, font_name="Microsoft YaHei"):
    if not content_list: return
    text_frame.clear()
    
    # 1. 获取文本框尺寸 (带更强的安全兜底)
    try:
        parent = text_frame._parent
        box_width_pt = parent.width.pt - Pt(10) 
        box_height_pt = parent.height.pt - Pt(10)
        
        # ⚠️ 关键修正：如果获取到的高度太小（比如小于 2英寸）
        # 强制认为它有一个标准正文框的高度 (约 5 英寸)
        if box_height_pt < Inches(2).pt:
            box_height_pt = Inches(5).pt
            
    except:
        # 完全获取不到时的默认值
        box_width_pt = Inches(8).pt
        box_height_pt = Inches(5).pt

    # 2. 定义字号列表 
    candidate_sizes = [32, 28, 24]
    best_size = 24

    # 3. 模拟排版
    for size in candidate_sizes:
        avg_char_width = size * 0.6 
        line_height = size * 1.2
        
        chars_per_line = max(1, int(box_width_pt / avg_char_width))
        
        total_lines = 0
        for line_text in content_list:
            text_len = len(str(line_text))
            if text_len == 0:
                total_lines += 1
                continue
            lines_needed = math.ceil(text_len / chars_per_line)
            total_lines += lines_needed
            
        estimated_height = total_lines * line_height
        
        # 如果能在高度限制内装下，就选用当前这个大字号
        if estimated_height <= box_height_pt:
            best_size = size
            break
    
    # 4. 应用字号
    for line in content_list:
        p = text_frame.add_paragraph()
        p.text = str(line)
        p.font.size = Pt(best_size) 
        p.font.name = font_name
        p.space_after = Pt(10)

def create_manual_table(slide, data, font_name="Microsoft YaHei"):
    """
    手动创建表格 (处理模板可能没有表格占位符的情况)
    """
    headers = data.headers
    rows = data.rows
    
    # === 动态计算表格高度 ===
    # 基础行高估算: 表头 0.5英寸 + 每行 0.4英寸
    row_count = len(rows)
    estimated_height = 0.5 + (row_count * 0.4)
    # 限制最大高度，防止画出幻灯片外面 (PPT一般高7.5英寸)
    final_height = min(estimated_height, 5.0) 
    
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(final_height) 

    # 创建表格形状
    shape = slide.shapes.add_table(row_count+1, len(headers), left, top, width, height)
    table = shape.table

    # 1. 填充表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192) # 经典蓝
        
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(18)      
            p.font.name = font_name   

    # 2. 填充数据行
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            if c_idx < len(headers):
                cell = table.cell(r_idx+1, c_idx)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(16)    
                    p.font.name = font_name 

# === 2. 布局配置 (来自 backend) ===
LAYOUT_CONFIG = {
    "academic": {
        "file": "templates/academic.pptx",
        "layouts": {
            "title_cover":  {"idx": 0, "title": 2, "sub": 3},
            "content_list": {"idx": 1, "title": 0, "body": 1},
            "chart":        {"idx": 1, "title": 0, "body": 1},
            "table":        {"idx": 1, "title": 0, "body": 1}, 
            "image_page":   {"idx": 1, "title": 0, "body": 1},
            "two_column":   {"idx": 2, "title": 0, "left": 1, "right": 2}
        }
    },
    "business": {
        "file": "templates/business.pptx",
        "layouts": {
            "title_cover":  {"idx": 0, "title": 0, "sub": 1},
            "content_list": {"idx": 1, "title": 0, "body": 1},
            "chart":        {"idx": 1, "title": 0, "body": 1},
            "table":        {"idx": 1, "title": 0, "body": 1}, 
            "image_page":   {"idx": 1, "title": 0, "body": 1},
            "two_column":   {"idx": 2, "title": 0, "left": 1, "right": 2}
        }
    },
    "teaching": {
        "file": "templates/teaching.pptx",
        "layouts": {
            "title_cover":  {"idx": 0, "title": 0, "sub": 13},
            "content_list": {"idx": 1, "title": 0, "body": 1},
            "chart":        {"idx": 1, "title": 0, "body": 1},
            "table":        {"idx": 1, "title": 0, "body": 1}, 
            "image_page":   {"idx": 1, "title": 0, "body": 1},
            "two_column":   {"idx": 2, "title": 0, "left": 1, "right": 2}
        }
    }
}

# === 3. 核心生成函数 ===
def create_pptx_file(data: PresentationData, theme: str = "academic") -> str:
    print(f"🎨 [Render] 开始渲染 PPT: {data.topic} (主题: {theme})")
    
    config = LAYOUT_CONFIG.get(theme, LAYOUT_CONFIG[theme])
    template_path = config["file"]
    
    if not os.path.exists(template_path):
        prs = Presentation() # 没有模板就用空白的
    else:
        prs = Presentation(template_path)

    layout_map = config["layouts"]
    
    # 定义全局字体，方便统一修改
    global_font = "Microsoft YaHei"

    for slide_data in data.slides:
        l_type = slide_data.layout
        print(f"   📄 处理页面 {slide_data.id}: {l_type}")

        # 1. 获取布局配置
        cfg = layout_map.get(l_type, layout_map["content_list"])
        slide_layout = prs.slide_layouts[cfg["idx"]]
        slide = prs.slides.add_slide(slide_layout)
        
        # 2. 填充通用标题
        try:
            if slide_data.title:
                slide.placeholders[cfg["title"]].text = slide_data.title
        except: pass
            
        # 3. 根据类型分发处理逻辑
        try:
            # --- Case A: 封面页 ---
            if l_type == "title_cover" and slide_data.subtitle:
                if slide_data.subtitle:
                    slide.placeholders[cfg["sub"]].text = slide_data.subtitle

            # --- Case B: 列表内容页 (使用 Auto-fit) ---
            elif l_type == "content_list":
                if slide_data.content and slide_data.content.bullet_points:
                    body_ph = slide.placeholders[cfg["body"]]
                    auto_fit_text(body_ph.text_frame, slide_data.content.bullet_points, font_name=global_font)
                elif slide_data.content and slide_data.content.text_body:
                    body_ph = slide.placeholders[cfg["body"]]
                    auto_fit_text(body_ph.text_frame, [slide_data.content.text_body], font_name=global_font)

            # --- Case C: 左右栏布局 ---
            elif l_type == "two_column":
                if slide_data.content:
                    if slide_data.content.content_left:
                        ph_left = slide.placeholders[cfg["left"]]
                        auto_fit_text(ph_left.text_frame, slide_data.content.content_left, font_name=global_font)
                    if slide_data.content.content_right:
                        ph_right = slide.placeholders[cfg["right"]]
                        auto_fit_text(ph_right.text_frame, slide_data.content.content_right, font_name=global_font)

            # --- Case D: 表格页 (新增) ---
            elif l_type == "table" and slide_data.table_data:
                # 如果有正文占位符，先清空或删除，防止遮挡
                if "body" in cfg and len(slide.placeholders) > cfg["body"]:
                    sp = slide.placeholders[cfg["body"]]
                    sp.element.getparent().remove(sp.element)
                
                create_manual_table(slide, slide_data.table_data, font_name=global_font)

            # --- Case E: 图表页 ---
            elif l_type == "chart" and slide_data.chart_data:
                chart_data = CategoryChartData()
                chart_data.categories = slide_data.chart_data.labels
                chart_data.add_series(slide_data.chart_data.title or "Series 1", slide_data.chart_data.values)
                
                # 尝试利用模板里的 Chart 占位符
                if "body" in cfg and len(slide.placeholders) > cfg["body"]:
                    ph = slide.placeholders[cfg["body"]]
                    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, ph.left, ph.top, ph.width, ph.height, chart_data)
                else:
                    # 默认位置
                    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2), Inches(8), Inches(4.5), chart_data)
                
                ph.element.getparent().remove(ph.element)


            # --- Case F: 图片处理 (通用) ---
            if slide_data.visual and slide_data.visual.need_image:
                prompt = slide_data.visual.image_prompt
                if prompt:
                    img_stream = get_image_stream(prompt)
                    if img_stream:
                        if l_type == "image_page":
                            # 大图居中
                            slide.shapes.add_picture(img_stream, Inches(1), Inches(2), width=Inches(8))
                            # 如果有 caption
                            if slide_data.visual.caption:
                                txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
                                p = txBox.text_frame.add_paragraph()
                                p.text = slide_data.visual.caption
                                p.font.size = Pt(12)
                                p.alignment = 2 # 居中
                            ph = slide.placeholders[cfg["body"]]
                            ph.element.getparent().remove(ph.element)
                        else:
                            # 装饰性小图 (右上角或右下角)
                            slide.shapes.add_picture(img_stream, Inches(6.5), Inches(5), width=Inches(3))
                            
        except Exception as e:
            print(f"⚠️ 页面 {slide_data.id} 渲染出错: {e}")
            continue

    # 保存
    os.makedirs("generated_ppts", exist_ok=True)
    filename = f"{uuid.uuid4()}.pptx"
    save_path = os.path.join("generated_ppts", filename)
    prs.save(save_path)
    print(f"✅ 文件已保存: {save_path}")
    
    return filename